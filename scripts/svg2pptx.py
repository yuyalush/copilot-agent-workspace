#!/usr/bin/env python3
"""
SVG → PPTX 変換スクリプト

SVG ファイルのシェイプ要素を PowerPoint のネイティブシェイプに変換し、
編集可能な .pptx ファイルを生成する。

Usage:
    python svg2pptx.py slides/                           # フォルダ名から自動命名
    python svg2pptx.py slides/ output.pptx               # 出力ファイル名を明示指定
    python svg2pptx.py slide01.svg slide02.svg out.pptx  # 個別ファイル指定
    python svg2pptx.py slide01.svg -o out.pptx --template base.pptx

Requirements:
    pip install python-pptx lxml
"""

import argparse
import math
import os
import re
import sys
import unicodedata
from pathlib import Path
from xml.etree import ElementTree as ET

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# SVG namespace
SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"
NSMAP = {"svg": SVG_NS}

# スライドサイズ (16:9 標準)
SLIDE_WIDTH_PX = 960
SLIDE_HEIGHT_PX = 540
SLIDE_WIDTH_EMU = Cm(25.4)   # 960px ≒ 25.4cm
SLIDE_HEIGHT_EMU = Cm(14.29)  # 540px ≒ 14.29cm

# px → EMU 変換係数
PX_TO_EMU = SLIDE_WIDTH_EMU / SLIDE_WIDTH_PX


def px(value: str | float) -> int:
    """px 値を EMU に変換"""
    if isinstance(value, str):
        value = float(value.replace("px", "").strip())
    return int(value * PX_TO_EMU)


def estimate_text_width_px(text: str, font_size_px: float) -> float:
    """テキストの幅を px 単位で推定（CJK 全角文字を考慮）"""
    width = 0.0
    for ch in text:
        # Unicode の East Asian Width で全角判定
        eaw = unicodedata.east_asian_width(ch)
        if eaw in ("F", "W"):
            # 全角文字: フォントサイズとほぼ同じ幅
            width += font_size_px * 0.95
        elif eaw == "A":
            # Ambiguous（丸数字等）: 少し広め
            width += font_size_px * 0.75
        else:
            # 半角（ASCII等）: フォントサイズの約 0.55 倍
            width += font_size_px * 0.55
    return width


def set_textbox_margins(txBox, left=0, top=0, right=0, bottom=0):
    """TextBox の内部マージンを EMU 単位で設定（デフォルトの 0.1 インチを除去）"""
    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("lIns", str(left))
        bodyPr.set("tIns", str(top))
        bodyPr.set("rIns", str(right))
        bodyPr.set("bIns", str(bottom))


def set_fill_opacity(shape, opacity: float):
    """シェイプの塗りつぶしに透明度を設定（OOXML の a:alpha を直接操作）

    opacity: 0.0（完全透明）〜 1.0（不透明）
    OOXML alpha val: 0 (透明) 〜 100000 (不透明)
    """
    if opacity >= 1.0:
        return  # 不透明なら何もしない

    alpha_val = int(opacity * 100000)

    # 塗りつぶしの透明度設定
    try:
        sp = shape._element
        spPr = sp.find(qn("a:spPr"), sp.nsmap) or sp.find(qn("p:spPr"), sp.nsmap)
        if spPr is None:
            # spPr が直下にない場合、全 namespace を探索
            for child in sp:
                if child.tag.endswith("spPr"):
                    spPr = child
                    break
        if spPr is None:
            return

        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is None:
            # fillRef 内を探す
            for desc in spPr.iter():
                if desc.tag.endswith("solidFill"):
                    solidFill = desc
                    break
        if solidFill is None:
            return

        # srgbClr を探す
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is None:
            return

        # 既存の alpha 要素を削除
        for existing in srgb.findall(qn("a:alpha")):
            srgb.remove(existing)

        # alpha 要素を追加（lxml SubElement を使用）
        alpha_ns = qn("a:alpha")
        alpha_elem = etree.SubElement(srgb, alpha_ns)
        alpha_elem.set("val", str(alpha_val))
    except Exception:
        pass  # 失敗しても続行


def set_line_opacity(shape, opacity: float):
    """シェイプの枠線に透明度を設定"""
    if opacity >= 1.0:
        return

    alpha_val = int(opacity * 100000)

    try:
        sp = shape._element
        # ln 要素を探す
        ln = None
        for desc in sp.iter():
            if desc.tag.endswith("}ln") or desc.tag == qn("a:ln"):
                ln = desc
                break
        if ln is None:
            return

        solidFill = ln.find(qn("a:solidFill"))
        if solidFill is None:
            return

        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is None:
            return

        for existing in srgb.findall(qn("a:alpha")):
            srgb.remove(existing)

        alpha_ns = qn("a:alpha")
        alpha_elem = etree.SubElement(srgb, alpha_ns)
        alpha_elem.set("val", str(alpha_val))
    except Exception:
        pass


def parse_color(color_str: str) -> RGBColor | None:
    """CSS カラー文字列を RGBColor に変換"""
    if not color_str or color_str == "none":
        return None

    color_str = color_str.strip()

    # #RRGGBB or #RGB
    if color_str.startswith("#"):
        hex_color = color_str.lstrip("#")
        if len(hex_color) == 3:
            hex_color = "".join(c * 2 for c in hex_color)
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    # rgb(r, g, b)
    match = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", color_str)
    if match:
        return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))

    # 名前付きカラー（主要なもの）
    named_colors = {
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "black": RGBColor(0x00, 0x00, 0x00),
        "red": RGBColor(0xFF, 0x00, 0x00),
        "blue": RGBColor(0x00, 0x00, 0xFF),
        "green": RGBColor(0x00, 0x80, 0x00),
        "gray": RGBColor(0x80, 0x80, 0x80),
        "grey": RGBColor(0x80, 0x80, 0x80),
    }
    return named_colors.get(color_str.lower())


def parse_gradient_ref(fill_str: str) -> str | None:
    """fill 値が url(#id) 参照かチェックし、グラデーション ID を返す"""
    if not fill_str:
        return None
    match = re.match(r"url\(#(.+?)\)", fill_str.strip())
    if match:
        return match.group(1)
    return None


def parse_gradient_defs(root) -> dict:
    """SVG の <defs> から linearGradient / radialGradient 定義を解析する

    Returns:
        dict: {gradient_id: {"type": "linear"|"radial", "stops": [...], ...}}
    """
    gradients = {}

    for elem in root.iter():
        tag = elem.tag
        if not isinstance(tag, str):
            continue
        if "}" in tag:
            tag = tag.split("}")[-1]

        if tag not in ("linearGradient", "radialGradient"):
            continue

        grad_id = elem.get("id")
        if not grad_id:
            continue

        # ストップ色を解析
        stops = []
        for child in elem:
            child_tag = child.tag
            if not isinstance(child_tag, str):
                continue
            if "}" in child_tag:
                child_tag = child_tag.split("}")[-1]
            if child_tag != "stop":
                continue

            offset = get_attr(child, "offset", "0")
            if offset.endswith("%"):
                offset = float(offset.rstrip("%")) / 100
            else:
                offset = float(offset)

            color = get_attr(child, "stop-color", "#000000")
            opacity = float(get_attr(child, "stop-opacity", "1"))
            stops.append({"offset": offset, "color": color, "opacity": opacity})

        if not stops:
            continue

        if tag == "linearGradient":
            x1 = float(elem.get("x1", "0").replace("%", "")) / (100 if "%" in elem.get("x1", "") else 1)
            y1 = float(elem.get("y1", "0").replace("%", "")) / (100 if "%" in elem.get("y1", "") else 1)
            x2 = float(elem.get("x2", "1").replace("%", "")) / (100 if "%" in elem.get("x2", "") else 1)
            y2 = float(elem.get("y2", "0").replace("%", "")) / (100 if "%" in elem.get("y2", "") else 1)

            angle_rad = math.atan2(y2 - y1, x2 - x1)
            angle_deg = math.degrees(angle_rad)
            if angle_deg < 0:
                angle_deg += 360

            gradients[grad_id] = {
                "type": "linear",
                "angle": angle_deg,
                "stops": stops,
            }

        elif tag == "radialGradient":
            cx = float(elem.get("cx", "0.5"))
            cy = float(elem.get("cy", "0.5"))
            r = float(elem.get("r", "0.5"))

            gradients[grad_id] = {
                "type": "radial",
                "cx": cx,
                "cy": cy,
                "r": r,
                "stops": stops,
            }

    return gradients


def apply_gradient_fill(shape, gradient_def: dict, opacity: float = 1.0):
    """PPTX シェイプにグラデーション塗りつぶしを適用する

    SVG の <linearGradient> / <radialGradient> を OOXML の a:gradFill に変換。
    """
    # python-pptx API でグラデーション塗りに切り替え
    shape.fill.gradient()

    # OOXML 要素を直接操作
    sp = shape._element
    spPr = None
    for child in sp:
        if child.tag.endswith("spPr"):
            spPr = child
            break
    if spPr is None:
        return

    gradFill = spPr.find(qn("a:gradFill"))
    if gradFill is None:
        return

    # グラデーションストップを再構築
    gsLst = gradFill.find(qn("a:gsLst"))
    if gsLst is None:
        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    else:
        for child in list(gsLst):
            gsLst.remove(child)

    for stop in gradient_def["stops"]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(int(stop["offset"] * 100000)))

        color = parse_color(stop["color"])
        if color:
            srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
            srgbClr.set("val", str(color))

            # stop-opacity と inherited opacity を合成
            combined_opacity = stop.get("opacity", 1.0) * opacity
            if combined_opacity < 1.0:
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", str(int(combined_opacity * 100000)))

    # 既存の方向指定を削除
    for child in list(gradFill):
        if child.tag == qn("a:lin") or child.tag == qn("a:path"):
            gradFill.remove(child)

    if gradient_def["type"] == "linear":
        lin = etree.SubElement(gradFill, qn("a:lin"))
        angle_ooxml = int(gradient_def["angle"] * 60000)
        lin.set("ang", str(angle_ooxml))
        lin.set("scaled", "0")

    elif gradient_def["type"] == "radial":
        path_elem = etree.SubElement(gradFill, qn("a:path"))
        path_elem.set("path", "circle")
        fillToRect = etree.SubElement(path_elem, qn("a:fillToRect"))
        cx_pct = int(gradient_def.get("cx", 0.5) * 100000)
        cy_pct = int(gradient_def.get("cy", 0.5) * 100000)
        fillToRect.set("l", str(cx_pct))
        fillToRect.set("t", str(cy_pct))
        fillToRect.set("r", str(100000 - cx_pct))
        fillToRect.set("b", str(100000 - cy_pct))


def parse_font_size(size_str: str) -> Pt:
    """フォントサイズ文字列を Pt に変換"""
    if not size_str:
        return Pt(14)
    # px → pt (概算: 1px ≒ 0.75pt)
    val = float(re.sub(r"[^\d.]", "", size_str))
    if "px" in size_str or size_str.replace(".", "").isdigit():
        val = val * 0.75
    return Pt(val)


def get_text_anchor(text_anchor: str) -> PP_ALIGN:
    """SVG text-anchor を PP_ALIGN にマッピング"""
    mapping = {
        "middle": PP_ALIGN.CENTER,
        "end": PP_ALIGN.RIGHT,
        "start": PP_ALIGN.LEFT,
    }
    return mapping.get(text_anchor, PP_ALIGN.LEFT)


def get_attr(elem, attr, default=None):
    """要素から属性を取得（style 属性内も検索）"""
    val = elem.get(attr)
    if val is not None:
        return val

    # style 属性から探す
    style = elem.get("style", "")
    if style:
        for part in style.split(";"):
            if ":" in part:
                key, value = part.split(":", 1)
                if key.strip() == attr:
                    return value.strip()

    return default


def _get_href(elem):
    """<a> 要素から href を取得する（href / xlink:href 両対応）"""
    href = elem.get("href") or elem.get(f"{{{XLINK_NS}}}href")
    return href


def _strip_tag(tag_str: str) -> str:
    """lxml のタグ文字列から namespace を除去して要素名だけ返す"""
    if not isinstance(tag_str, str):
        return ""
    if "}" in tag_str:
        return tag_str.split("}")[-1]
    return tag_str


def parse_translate(elem) -> tuple[float, float]:
    """transform="translate(x, y)" から (dx, dy) を抽出する"""
    transform = elem.get("transform", "")
    # translate(x, y) — カンマ/スペース区切りの2引数
    m = re.search(r'translate\(\s*([\d.eE+-]+)[\s,]+([\d.eE+-]+)\s*\)', transform)
    if m:
        return float(m.group(1)), float(m.group(2))
    # translate(x) — y 省略時は 0
    m = re.search(r'translate\(\s*([\d.eE+-]+)\s*\)', transform)
    if m:
        return float(m.group(1)), 0.0
    return 0.0, 0.0


def collect_text_content(elem, href: str | None = None):
    """<text> 要素の直接テキストと <tspan> / <a> 子要素のテキストを収集"""
    texts = []

    # 直接テキスト
    if elem.text and elem.text.strip():
        texts.append({
            "text": elem.text.strip(),
            "font_size": get_attr(elem, "font-size"),
            "font_weight": get_attr(elem, "font-weight"),
            "fill": get_attr(elem, "fill"),
            "font_family": get_attr(elem, "font-family"),
            "href": href,
        })

    # 子要素（<tspan> / <a>）
    for child in elem:
        tag = _strip_tag(child.tag)

        if tag == "tspan" and child.text and child.text.strip():
            texts.append({
                "text": child.text.strip(),
                "font_size": get_attr(child, "font-size") or get_attr(elem, "font-size"),
                "font_weight": get_attr(child, "font-weight") or get_attr(elem, "font-weight"),
                "fill": get_attr(child, "fill") or get_attr(elem, "fill"),
                "font_family": get_attr(child, "font-family") or get_attr(elem, "font-family"),
                "href": href,
            })

        elif tag == "a":
            # <a href="..."> 内の <tspan> やテキストを再帰的に収集
            link_href = _get_href(child) or href
            texts.extend(collect_text_content(child, href=link_href))

    return texts


def apply_stroke_dasharray(shape, elem):
    """stroke-dasharray を PPTX のダッシュスタイルに変換"""
    dasharray = get_attr(elem, "stroke-dasharray")
    if not dasharray or dasharray == "none":
        return
    try:
        ln = None
        sp = shape._element
        for desc in sp.iter():
            if desc.tag.endswith("}ln") or desc.tag == qn("a:ln"):
                ln = desc
                break
        if ln is None:
            return
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    except Exception:
        pass


def remove_shadow(shape):
    """シェイプからデフォルトの影（ドロップシャドウ）を除去する

    python-pptx が追加するシェイプにはデフォルトで影が付くことがある。
    OOXML では a:effectLst を空にすることで影を無効化できる。
    """
    try:
        sp = shape._element
        # spPr を探す
        spPr = None
        for child in sp:
            if child.tag.endswith("spPr"):
                spPr = child
                break
        if spPr is None:
            return

        # 既存の effectLst を削除
        for eff in list(spPr.findall(qn("a:effectLst"))):
            spPr.remove(eff)

        # 空の effectLst を追加（"効果なし" を明示）
        etree.SubElement(spPr, qn("a:effectLst"))
    except Exception:
        pass


def get_effective_opacity(elem, inherited_opacity: float = 1.0) -> float:
    """要素自身の opacity と親から継承した opacity を合成して返す"""
    own_opacity = get_attr(elem, "opacity")
    if own_opacity is not None:
        return float(own_opacity) * inherited_opacity
    return inherited_opacity


def set_rounded_rect_radius(shape, rx_px: float, width_px: float, height_px: float):
    """RoundedRectangle の角丸半径を SVG の rx 値に正確に合わせる

    OOXML の roundRect の adj ガイド:
        角丸半径 = adj * min(width, height) / 100000
    したがって:
        adj = rx / min(width, height) * 100000

    デフォルト adj=16667（≒1/6）は多くの場合大きすぎるため、
    SVG で指定された rx に合わせて補正する。
    """
    min_dim = min(width_px, height_px)
    if min_dim <= 0:
        return
    adj_val = int((rx_px / min_dim) * 100000)
    adj_val = max(0, min(adj_val, 50000))  # 0〜50000 にクランプ

    try:
        sp = shape._element
        # spPr を探す（p:spPr or a:spPr）
        spPr = None
        for child in sp:
            if child.tag.endswith("spPr"):
                spPr = child
                break
        if spPr is None:
            return

        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is None:
            return

        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        else:
            # 既存の adj を削除して上書き
            for gd in list(avLst.findall(qn("a:gd"))):
                avLst.remove(gd)

        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {adj_val}")
    except Exception:
        pass


def add_rect(slide, elem, inherited_opacity: float = 1.0, gradients: dict | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """<rect> → AutoShape (Rectangle or Rounded Rectangle)"""
    x = px(float(get_attr(elem, "x", "0")) + offset_x)
    y = px(float(get_attr(elem, "y", "0")) + offset_y)
    w = px(get_attr(elem, "width", "0"))
    h = px(get_attr(elem, "height", "0"))
    rx = get_attr(elem, "rx")
    rx_val = float(rx) if rx else 0

    if rx_val > 0:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    remove_shadow(shape)

    # 角丸半径を SVG の rx に合わせる（デフォルトだと大きすぎる）
    if rx_val > 0:
        w_px = float(get_attr(elem, "width", "0"))
        h_px = float(get_attr(elem, "height", "0"))
        set_rounded_rect_radius(shape, rx_val, w_px, h_px)

    # 実効 opacity を計算
    effective_opacity = get_effective_opacity(elem, inherited_opacity)

    # 塾りつぶし（グラデーション対応）
    fill_str = get_attr(elem, "fill", "#FFFFFF")
    grad_id = parse_gradient_ref(fill_str)
    if grad_id and gradients and grad_id in gradients:
        apply_gradient_fill(shape, gradients[grad_id], effective_opacity)
    else:
        fill_color = parse_color(fill_str)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            if effective_opacity < 1.0:
                set_fill_opacity(shape, effective_opacity)
        else:
            shape.fill.background()

    # 枠線
    stroke = get_attr(elem, "stroke")
    if stroke and stroke != "none":
        stroke_color = parse_color(stroke)
        if stroke_color:
            shape.line.color.rgb = stroke_color
            stroke_width = get_attr(elem, "stroke-width", "1")
            shape.line.width = Pt(float(stroke_width))
            if effective_opacity < 1.0:
                set_line_opacity(shape, effective_opacity)
        # ダッシュスタイル
        apply_stroke_dasharray(shape, elem)
    else:
        # SVG で stroke 未指定の場合、枠線を完全に消す
        shape.line.fill.background()
        shape.line.width = Emu(0)

    return shape


def add_text(slide, elem, href: str | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """<text> → TextBox（改良版: CJK幅推定・ベースライン補正・マージン除去・ハイパーリンク対応）"""
    x_val = float(get_attr(elem, "x", "0")) + offset_x
    y_val = float(get_attr(elem, "y", "0")) + offset_y
    text_anchor = get_attr(elem, "text-anchor", "start")

    texts = collect_text_content(elem, href=href)
    if not texts:
        return None

    # フォントサイズ取得（px 単位）
    font_size_str = get_attr(elem, "font-size", "14")
    font_size_val = float(re.sub(r"[^\d.]", "", font_size_str))

    # --- テキスト幅を CJK 対応で推定 ---
    full_text = " ".join(t["text"] for t in texts)
    estimated_width = estimate_text_width_px(full_text, font_size_val)
    # 余裕を少し追加（フォントによるばらつき対策）
    estimated_width = max(estimated_width * 1.05, font_size_val * 2)

    # --- テキスト高さ推定 ---
    # tspan が複数ある場合、各行の高さを個別に計算
    line_heights = []
    for t in texts:
        t_font_str = t.get("font_size") or font_size_str
        t_font_val = float(re.sub(r"[^\d.]", "", t_font_str))
        line_heights.append(t_font_val * 1.4)  # 行間 1.4 倍
    estimated_height = sum(line_heights) if line_heights else font_size_val * 1.4

    # --- text-anchor に応じて x 座標を調整 ---
    if text_anchor == "middle":
        x_val -= estimated_width / 2
    elif text_anchor == "end":
        x_val -= estimated_width

    # --- SVG ベースライン → PPTX 左上原点 補正 ---
    # SVG の y はテキストのベースライン（文字の底辺付近）
    # PPTX TextBox は左上原点なので、フォントのアセント分を引く
    # アセント ≒ フォントサイズの 80% が一般的
    y_val -= font_size_val * 0.85

    txBox = slide.shapes.add_textbox(
        px(x_val), px(y_val),
        px(estimated_width), px(estimated_height)
    )

    # 内部マージンを 0 にする（デフォルトの 91440 EMU = 0.1'' を除去）
    set_textbox_margins(txBox, left=0, top=0, right=0, bottom=0)

    tf = txBox.text_frame
    tf.word_wrap = False  # 1行テキストでは折り返しを無効化

    # 複数行ある場合のみ word_wrap を有効化
    if len(texts) > 1:
        tf.word_wrap = True

    for i, text_info in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = get_text_anchor(text_anchor)

        # 段落間スペースを 0 にする（テキストの重なり防止）
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        run = p.add_run()
        run.text = text_info["text"]

        # フォント設定
        if text_info.get("font_size"):
            run.font.size = parse_font_size(text_info["font_size"])
        if text_info.get("font_weight") == "bold":
            run.font.bold = True
        if text_info.get("fill"):
            color = parse_color(text_info["fill"])
            if color:
                run.font.color.rgb = color
        if text_info.get("font_family"):
            font_name = text_info["font_family"].split(",")[0].strip().strip("'\"")
            run.font.name = font_name

        # ハイパーリンク
        if text_info.get("href"):
            run.hyperlink.address = text_info["href"]

    return txBox


def add_circle(slide, elem, inherited_opacity: float = 1.0, gradients: dict | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """<circle> → Oval AutoShape"""
    cx = float(get_attr(elem, "cx", "0")) + offset_x
    cy = float(get_attr(elem, "cy", "0")) + offset_y
    r = float(get_attr(elem, "r", "0"))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px(cx - r), px(cy - r),
        px(r * 2), px(r * 2)
    )
    remove_shadow(shape)

    effective_opacity = get_effective_opacity(elem, inherited_opacity)

    fill_str = get_attr(elem, "fill")
    grad_id = parse_gradient_ref(fill_str)
    if grad_id and gradients and grad_id in gradients:
        apply_gradient_fill(shape, gradients[grad_id], effective_opacity)
    else:
        fill_color = parse_color(fill_str)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            if effective_opacity < 1.0:
                set_fill_opacity(shape, effective_opacity)
        else:
            shape.fill.background()

    stroke = get_attr(elem, "stroke")
    if stroke and stroke != "none":
        stroke_color = parse_color(stroke)
        if stroke_color:
            shape.line.color.rgb = stroke_color
            if effective_opacity < 1.0:
                set_line_opacity(shape, effective_opacity)
    else:
        shape.line.fill.background()
        shape.line.width = Emu(0)

    return shape


def add_ellipse(slide, elem, inherited_opacity: float = 1.0, gradients: dict | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """<ellipse> → Oval AutoShape"""
    cx = float(get_attr(elem, "cx", "0")) + offset_x
    cy = float(get_attr(elem, "cy", "0")) + offset_y
    rx = float(get_attr(elem, "rx", "0"))
    ry = float(get_attr(elem, "ry", "0"))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px(cx - rx), px(cy - ry),
        px(rx * 2), px(ry * 2)
    )
    remove_shadow(shape)

    effective_opacity = get_effective_opacity(elem, inherited_opacity)

    fill_str = get_attr(elem, "fill")
    grad_id = parse_gradient_ref(fill_str)
    if grad_id and gradients and grad_id in gradients:
        apply_gradient_fill(shape, gradients[grad_id], effective_opacity)
    else:
        fill_color = parse_color(fill_str)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            if effective_opacity < 1.0:
                set_fill_opacity(shape, effective_opacity)

    return shape


def add_line(slide, elem, inherited_opacity: float = 1.0, offset_x: float = 0.0, offset_y: float = 0.0):
    """<line> → Connector"""
    x1 = px(float(get_attr(elem, "x1", "0")) + offset_x)
    y1 = px(float(get_attr(elem, "y1", "0")) + offset_y)
    x2 = px(float(get_attr(elem, "x2", "0")) + offset_x)
    y2 = px(float(get_attr(elem, "y2", "0")) + offset_y)

    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        x1, y1, x2, y2
    )

    effective_opacity = get_effective_opacity(elem, inherited_opacity)

    stroke = get_attr(elem, "stroke", "#000000")
    stroke_color = parse_color(stroke)
    if stroke_color:
        connector.line.color.rgb = stroke_color
        if effective_opacity < 1.0:
            set_line_opacity(connector, effective_opacity)

    stroke_width = get_attr(elem, "stroke-width", "1")
    connector.line.width = Pt(float(stroke_width))

    # ダッシュスタイル
    apply_stroke_dasharray(connector, elem)

    return connector


def add_polygon(slide, elem, inherited_opacity: float = 1.0, gradients: dict | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """<polygon> → Freeform shape"""
    points_str = get_attr(elem, "points", "")
    if not points_str:
        return None

    # "x1,y1 x2,y2 x3,y3" をパース
    point_pairs = points_str.strip().split()
    points = []
    for pair in point_pairs:
        parts = pair.split(",")
        if len(parts) == 2:
            points.append((float(parts[0]) + offset_x, float(parts[1]) + offset_y))

    if len(points) < 3:
        return None

    # FreeformBuilder で多角形を作成（絶対座標を使用）
    builder = slide.shapes.build_freeform(
        px(points[0][0]), px(points[0][1])
    )
    # 残りの点へ線を引く（絶対座標）
    segments = []
    for point in points[1:]:
        segments.append((px(point[0]), px(point[1])))
    # 閉じる（最初の点に戻る）
    segments.append((px(points[0][0]), px(points[0][1])))
    builder.add_line_segments(segments)

    shape = builder.convert_to_shape()
    remove_shadow(shape)

    effective_opacity = get_effective_opacity(elem, inherited_opacity)

    fill_str = get_attr(elem, "fill")
    grad_id = parse_gradient_ref(fill_str)
    if grad_id and gradients and grad_id in gradients:
        apply_gradient_fill(shape, gradients[grad_id], effective_opacity)
    else:
        fill_color = parse_color(fill_str)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            if effective_opacity < 1.0:
                set_fill_opacity(shape, effective_opacity)

    stroke = get_attr(elem, "stroke")
    if stroke and stroke != "none":
        stroke_color = parse_color(stroke)
        if stroke_color:
            shape.line.color.rgb = stroke_color
    else:
        shape.line.fill.background()
        shape.line.width = Emu(0)

    return shape


def process_group(slide, group_elem, inherited_opacity: float = 1.0, gradients: dict | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """<g> 内の要素を再帰的に処理（グループの opacity・translate を子要素に継承）"""
    group_opacity = get_effective_opacity(group_elem, inherited_opacity)
    # transform="translate(x, y)" のオフセットを累積
    dx, dy = parse_translate(group_elem)
    new_offset_x = offset_x + dx
    new_offset_y = offset_y + dy
    for child in group_elem:
        process_element(slide, child, group_opacity, gradients, new_offset_x, new_offset_y)


def process_element(slide, elem, inherited_opacity: float = 1.0, gradients: dict | None = None, offset_x: float = 0.0, offset_y: float = 0.0):
    """SVG 要素を対応する PPTX シェイプに変換"""
    tag = elem.tag
    # lxml ではコメント等で tag が callable になる場合がある
    if not isinstance(tag, str):
        return
    # namespace を除去
    if "}" in tag:
        tag = tag.split("}")[-1]

    # defs はスキップ（グラデーション定義は事前に解析済み）
    if tag == "defs":
        return

    # グラデーション対応ハンドラ
    gradient_handlers = {
        "rect": add_rect,
        "circle": add_circle,
        "ellipse": add_ellipse,
        "polygon": add_polygon,
        "g": process_group,
    }

    # line は gradients 不要
    if tag == "line":
        try:
            add_line(slide, elem, inherited_opacity, offset_x, offset_y)
        except Exception as e:
            print(f"  ⚠ {tag} 要素の変換をスキップ: {e}", file=sys.stderr)
        return

    # テキストは opacity 未対応（テキスト自体の透過は稀）
    if tag == "text":
        try:
            add_text(slide, elem, offset_x=offset_x, offset_y=offset_y)
        except Exception as e:
            print(f"  ⚠ {tag} 要素の変換をスキップ: {e}", file=sys.stderr)
        return

    # <a> リンク: 子要素を href 付きで処理
    if tag == "a":
        href = _get_href(elem)
        for child in elem:
            child_tag = _strip_tag(child.tag)
            if child_tag == "text":
                try:
                    # collect_text_content に href を伝播
                    add_text(slide, child, href=href, offset_x=offset_x, offset_y=offset_y)
                except Exception as e:
                    print(f"  ⚠ a>text 要素の変換をスキップ: {e}", file=sys.stderr)
            else:
                process_element(slide, child, inherited_opacity, gradients, offset_x, offset_y)
        return

    handler = gradient_handlers.get(tag)
    if handler:
        try:
            handler(slide, elem, inherited_opacity, gradients, offset_x, offset_y)
        except Exception as e:
            print(f"  ⚠ {tag} 要素の変換をスキップ: {e}", file=sys.stderr)


def convert_svg_to_slide(prs: Presentation, svg_path: str):
    """1つの SVG ファイルを1枚のスライドに変換"""
    print(f"  📄 {Path(svg_path).name}")

    tree = etree.parse(svg_path)
    root = tree.getroot()

    # 空白レイアウトのスライドを追加
    blank_layout = prs.slide_layouts[6]  # 通常 index 6 が空白
    slide = prs.slides.add_slide(blank_layout)

    # グラデーション定義を事前に解析
    gradients = parse_gradient_defs(root)

    # SVG の子要素を順番に処理
    for elem in root:
        process_element(slide, elem, gradients=gradients)

    return slide


def find_svg_files(paths: list[str]) -> list[str]:
    """パスリストから SVG ファイルを収集（ディレクトリの場合は中身を探索）"""
    svg_files = []
    for path_str in paths:
        p = Path(path_str)
        if p.is_dir():
            svg_files.extend(sorted(p.glob("*.svg")))
        elif p.suffix.lower() == ".svg" and p.exists():
            svg_files.append(p)
        else:
            print(f"  ⚠ スキップ: {path_str}", file=sys.stderr)
    return [str(f) for f in svg_files]


def _auto_pptx_name(inputs: list[str]) -> str:
    """入力パスからPPTXファイル名を自動生成する。

    フォルダ名が ``YYYY-MM-DD_HHmmss_テーマ名`` 形式の場合、
    ``<入力フォルダ>/YYYY-MM-DD_テーマ名.pptx`` を返す（時分秒を省略）。
    それ以外はフォルダ/ファイル名をそのまま使い、入力フォルダ内に保存する。
    """
    # 入力ディレクトリを特定
    p = Path(inputs[0]).resolve()
    if p.is_dir() or not p.suffix:
        # ディレクトリ指定（存在する or 拡張子なし = ディレクトリ風パス）
        dir_path = p
        dir_name = p.name
    else:
        # ファイル指定 → 親ディレクトリ
        dir_path = p.parent
        dir_name = p.parent.name

    # YYYY-MM-DD_HHmmss_テーマ名 パターンにマッチ → 時分秒を省略
    m = re.match(r"^(\d{4}-\d{2}-\d{2})_\d{6}_(.+)$", dir_name)
    if m:
        pptx_name = f"{m.group(1)}_{m.group(2)}.pptx"
        return str(dir_path / pptx_name)

    # YYYY-MM-DD_テーマ名（HHmmss なし）パターン
    m2 = re.match(r"^(\d{4}-\d{2}-\d{2})_(.+)$", dir_name)
    if m2:
        return str(dir_path / f"{dir_name}.pptx")

    # パターン不一致 → フォルダ名をそのまま使う
    if dir_name:
        return str(dir_path / f"{dir_name}.pptx")

    return "output.pptx"


def main():
    parser = argparse.ArgumentParser(
        description="SVG スライドを編集可能な PPTX に変換する"
    )
    parser.add_argument(
        "inputs", nargs="+",
        help="SVG ファイルまたはディレクトリのパス"
    )
    parser.add_argument(
        "-o", "--output", default=None,
        help="出力 PPTX ファイル名 (省略時はフォルダ名から自動生成)"
    )
    parser.add_argument(
        "--template",
        help="ベースとなる PPTX テンプレートファイル"
    )

    args = parser.parse_args()

    # 最後の引数が .pptx の場合は出力ファイル名として扱う
    if args.inputs[-1].endswith(".pptx"):
        args.output = args.inputs.pop()

    # 出力ファイル名の自動生成
    if args.output is None:
        args.output = _auto_pptx_name(args.inputs)

    # SVG ファイル収集
    svg_files = find_svg_files(args.inputs)
    if not svg_files:
        print("❌ SVG ファイルが見つかりません", file=sys.stderr)
        sys.exit(1)

    print(f"🔄 {len(svg_files)} 枚のスライドを変換中...")

    # Presentation 作成
    if args.template and Path(args.template).exists():
        prs = Presentation(args.template)
        print(f"  📎 テンプレート: {args.template}")
    else:
        prs = Presentation()

    # スライドサイズ設定 (16:9)
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    # 各 SVG をスライドに変換
    for svg_file in svg_files:
        convert_svg_to_slide(prs, svg_file)

    # プロパティをクリア（デフォルトテンプレートの作者情報を除去）
    prs.core_properties.author = ""
    prs.core_properties.last_modified_by = ""
    prs.core_properties.comments = ""

    # 保存
    prs.save(args.output)
    print(f"✅ 保存完了: {args.output}")


if __name__ == "__main__":
    main()
